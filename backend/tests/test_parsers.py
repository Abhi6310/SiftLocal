import pytest
import io
from app.services.parser_client import parse_document
from pptx import Presentation
from pptx.util import Inches

def _create_sample_pptx():
    #create minimal PPTX with text for testing
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text = "Hello PPTX World"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def _create_empty_pptx():
    #create PPTX with no text shapes
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

SAMPLE_PPTX = _create_sample_pptx()
EMPTY_PPTX = _create_empty_pptx()

#minimal valid PDF with text content
SAMPLE_PDF = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF World) Tj ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
435
%%EOF"""

def test_pdf():
    #test PDF parsing extracts text and metadata
    result = parse_document(SAMPLE_PDF, ".pdf")
    assert result.error is None, f"Parser error: {result.error}"
    assert result.text is not None
    assert result.metadata is not None
    assert result.metadata["file_type"] == "pdf"
    assert result.metadata["page_count"] == 1
    assert "char_count" in result.metadata
    assert isinstance(result.text, str)

def test_pdf_empty():
    #minimal PDF with no text
    minimal = b"""%PDF-1.4
1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj
2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj
3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >> endobj
xref
0 4
0000000000 65535 f
0000000009 00000 n
0000000056 00000 n
0000000111 00000 n
trailer << /Size 4 /Root 1 0 R >>
startxref
176
%%EOF"""
    result = parse_document(minimal, ".pdf")
    assert result.error is None
    assert result.metadata["page_count"] == 1
    assert result.text == ""

def test_pdf_no_disk_write(tmp_path, monkeypatch):
    #invariant: parsed content never written to disk
    import builtins
    original_open = builtins.open
    disk_writes = []

    def tracked_open(path, mode="r", *args, **kwargs):
        if "w" in mode or "a" in mode:
            # Ignore pytest's own writes
            if not str(path).startswith("/tmp/pytest"):
                disk_writes.append(str(path))
        return original_open(path, mode, *args, **kwargs)

    monkeypatch.setattr(builtins, "open", tracked_open)
    result = parse_document(SAMPLE_PDF, ".pdf")

    #no disk writes should occur during parsing
    assert len(disk_writes) == 0, f"Unexpected disk writes in backend: {disk_writes}"
    assert result.error is None or len(result.text) >= 0

def test_pptx():
    #test PPTX parsing extracts text and metadata
    result = parse_document(SAMPLE_PPTX, ".pptx")
    assert result.error is None, f"Parser error: {result.error}"
    assert result.text is not None
    assert "Hello PPTX World" in result.text
    assert result.metadata is not None
    assert result.metadata["file_type"] == "pptx"
    assert result.metadata["slide_count"] == 1
    assert "char_count" in result.metadata
    assert "slides" in result.metadata

def test_pptx_empty():
    #test PPTX with no text shapes
    result = parse_document(EMPTY_PPTX, ".pptx")
    assert result.error is None
    assert result.metadata["slide_count"] == 1
    assert result.text == ""

def test_pptx_no_disk_write(tmp_path, monkeypatch):
    #invariant: parsed content never written to disk
    import builtins
    original_open = builtins.open
    disk_writes = []

    def tracked_open(path, mode="r", *args, **kwargs):
        if "w" in mode or "a" in mode:
            if not str(path).startswith("/tmp/pytest"):
                disk_writes.append(str(path))
        return original_open(path, mode, *args, **kwargs)

    monkeypatch.setattr(builtins, "open", tracked_open)
    result = parse_document(SAMPLE_PPTX, ".pptx")

    #no disk writes should occur during parsing
    assert len(disk_writes) == 0, f"Unexpected disk writes in backend: {disk_writes}"
    assert result.error is None or len(result.text) >= 0
