import io
from pptx import Presentation

def parse_pptx(content: bytes) -> dict:
    #extract text and structure from PPTX bytes
    prs = Presentation(io.BytesIO(content))
    slides_text = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content.append(shape.text)
        slides_text.append({
            "slide_number": slide_idx,
            "text": "\n".join(slide_content)
        })
    full_text = "\n\n".join([s["text"] for s in slides_text if s["text"]])
    return {
        "text": full_text,
        "metadata": {
            "file_type": "pptx",
            "slide_count": len(prs.slides),
            "char_count": len(full_text),
            "slides": [{"slide_number": s["slide_number"], "char_count": len(s["text"])} for s in slides_text]
        }
    }
